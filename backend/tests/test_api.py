import os
import io
import sys
import pytest

# Ensure backend root is on sys.path for tests
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from fastapi.testclient import TestClient
from reportlab.pdfgen import canvas
from main import app

client = TestClient(app)

@pytest.fixture
def dummy_pdf():
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer)
    c.drawString(100, 100, "Hello World")
    c.showPage()
    c.drawString(100, 100, "Page 2")
    c.showPage()
    c.save()
    buffer.seek(0)
    return buffer.read()

def test_read_root():
    response = client.get("/")
    assert response.status_code == 200
    assert response.json() == {"message": "PDF Converter API is running"}

def test_merge_pdfs(dummy_pdf):
    files = [
        ('files', ('test1.pdf', dummy_pdf, 'application/pdf')),
        ('files', ('test2.pdf', dummy_pdf, 'application/pdf'))
    ]
    response = client.post("/api/process/merge", files=files)
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"

def test_split_pdf(dummy_pdf):
    files = {'file': ('test.pdf', dummy_pdf, 'application/pdf')}
    data = {'page_range': '1-1'}
    response = client.post("/api/process/split", files=files, data=data)
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"

def test_organize_pdf(dummy_pdf):
    files = {'file': ('test.pdf', dummy_pdf, 'application/pdf')}
    data = {'page_order': '1,0', 'rotation': '{}'}
    response = client.post("/api/process/organize", files=files, data=data)
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"  # organize returns pdf without explicit media_type in code but FileResponse often infers? 
    # Actually my organize code didn't set media_type, so it might default or be missing. Check implementation.
    # It sets filename=.pdf, so FastAPI might infer.

def test_lock_pdf(dummy_pdf):
    files = {'file': ('test.pdf', dummy_pdf, 'application/pdf')}
    data = {'password': 'secure'}
    response = client.post("/api/process/lock", files=files, data=data)
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"

@pytest.fixture
def dummy_docx():
    import docx
    buffer = io.BytesIO()
    doc = docx.Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("This is a sample paragraph for testing Word to PDF conversion.")
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()

@pytest.fixture
def dummy_pptx():
    from pptx import Presentation
    buffer = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Test Presentation"
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()

@pytest.fixture
def dummy_xlsx():
    import pandas as pd
    buffer = io.BytesIO()
    df = pd.DataFrame({'Name': ['Alice', 'Bob'], 'Age': [25, 30]})
    with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
        df.to_excel(writer, sheet_name='Sheet1', index=False)
    buffer.seek(0)
    return buffer.read()

def test_word_to_pdf(dummy_docx):
    files = {'file': ('test.docx', dummy_docx, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')}
    response = client.post("/api/process/word-to-pdf", files=files)
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"

def test_ppt_to_pdf(dummy_pptx):
    files = {'file': ('test.pptx', dummy_pptx, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
    response = client.post("/api/process/ppt-to-pdf", files=files)
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"

def test_excel_to_pdf(dummy_xlsx):
    files = {'file': ('test.xlsx', dummy_xlsx, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')}
    response = client.post("/api/process/excel-to-pdf", files=files)
    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"

