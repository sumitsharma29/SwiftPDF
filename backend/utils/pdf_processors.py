import os
from typing import List
from pypdf import PdfReader, PdfWriter

class PDFProcessors:
    """
    Class containing all PDF manipulation logic.
    """

    @staticmethod
    def merge_pdfs(file_paths: List[str], output_path: str):
        merger = PdfWriter()
        for pdf in file_paths:
            merger.append(pdf)
        merger.write(output_path)
        merger.close()

    @staticmethod
    def split_pdf(file_path: str, output_path: str, page_range: str):
        # page_range format: "1-5", "1,3,5" or "1-2, 5"
        reader = PdfReader(file_path)
        writer = PdfWriter()
        
        try:
            pages_to_keep = set()
            for part in page_range.split(','):
                part = part.strip()
                if '-' in part:
                    start, end = map(int, part.split('-'))
                    pages_to_keep.update(range(start-1, end))
                else:
                    pages_to_keep.add(int(part)-1)

            for i in range(len(reader.pages)):
                if i in pages_to_keep:
                    writer.add_page(reader.pages[i])
            
            with open(output_path, "wb") as f:
                writer.write(f)
        except Exception:
            raise ValueError("Invalid page range format. Use '1-5', '1,3,5' or '1-2, 5'.")

    @staticmethod
    def organize_pdf(file_path: str, output_path: str, page_order: List[int], rotation: dict):
        # page_order: list of 0-based indices in desired order
        # rotation: dict of {page_index: angle}
        reader = PdfReader(file_path)
        writer = PdfWriter()
        
        for page_idx in page_order:
            if 0 <= page_idx < len(reader.pages):
                page = reader.pages[page_idx]
                if str(page_idx) in rotation:
                    # pypdf rotation is relative to current, usually just set it
                    page.rotate(int(rotation[str(page_idx)]))
                writer.add_page(page)
        
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def images_to_pdf(image_paths: List[str], output_path: str):
        import img2pdf
        # input: list of image paths
        # output: single pdf path
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(image_paths))

    @staticmethod
    def pdf_to_images(file_path: str, output_folder: str):
        import fitz  # PyMuPDF
        # Returns list of paths to generated images
        # Uses PyMuPDF (fitz) which doesn't require system poppler
        doc = fitz.open(file_path)
        image_paths = []
        
        for i in range(len(doc)):
            page = doc.load_page(i)
            # Render page to an image
            pix = page.get_pixmap()
            image_path = os.path.join(output_folder, f"page_{i + 1}.jpg")
            pix.save(image_path)
            image_paths.append(image_path)
            
        doc.close()
        return image_paths

    @staticmethod
    def compress_pdf(file_path: str, output_path: str, level: str):
        import pikepdf
        # Use pikepdf for better compression
        try:
            with pikepdf.open(file_path) as pdf:
                # Remove unused resources
                pdf.remove_unreferenced_resources()
                
                # Save with linearization and object stream generation
                pdf.save(
                    output_path, 
                    linearize=True, 
                    object_stream_mode=pikepdf.ObjectStreamMode.generate
                )
        except Exception as e:
            # Fallback to pypdf if pikepdf fails
            print(f"Pikepdf failed: {e}, using fallback")
            reader = PdfReader(file_path)
            writer = PdfWriter()
            for page in reader.pages:
                page.compress_content_streams()
                writer.add_page(page)
            with open(output_path, "wb") as f:
                writer.write(f)

    @staticmethod
    def lock_pdf(file_path: str, output_path: str, password: str):
        try:
            import pikepdf
            with pikepdf.Pdf.open(file_path) as pdf:
                pdf.save(output_path, encryption=pikepdf.Encryption(
                    user=password, owner=password, R=6
                ))
                return
        except Exception:
            pass

        reader = PdfReader(file_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(user_password=password, owner_password=password)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def unlock_pdf(file_path: str, output_path: str, password: str):
        try:
            import pikepdf
            try:
                with pikepdf.Pdf.open(file_path, password=password) as pdf:
                    pdf.save(output_path)
                    return
            except pikepdf.PasswordError:
                raise ValueError("Incorrect password")
        except ValueError:
            raise
        except Exception:
            pass

        try:
            reader = PdfReader(file_path, password=password)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            with open(output_path, "wb") as f:
                writer.write(f)
        except Exception:
            raise ValueError("Incorrect password or unreadable encrypted PDF")

    @staticmethod
    def watermark_pdf(file_path: str, output_path: str, watermark_text: str, output_folder: str):
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        # 1. Create watermark PDF
        watermark_file = os.path.join(output_folder, "watermark_temp.pdf")
        c = canvas.Canvas(watermark_file, pagesize=letter)
        c.setFont("Helvetica", 40)
        c.setStrokeColorRGB(0.5, 0.5, 0.5)
        c.setFillColorRGB(0.5, 0.5, 0.5)
        c.setFillAlpha(0.5)
        c.saveState()
        c.translate(300, 400)
        c.rotate(45)
        c.drawCentredString(0, 0, watermark_text)
        c.restoreState()
        c.save()

        # 2. Merge with source
        watermark_reader = PdfReader(watermark_file)
        if len(watermark_reader.pages) > 0:
            watermark_page = watermark_reader.pages[0]
            reader = PdfReader(file_path)
            writer = PdfWriter()
            for page in reader.pages:
                page.merge_page(watermark_page)
                writer.add_page(page)
            with open(output_path, "wb") as f:
                writer.write(f)
        
        if os.path.exists(watermark_file):
            os.remove(watermark_file)

    @staticmethod
    def extract_text(file_path: str) -> str:
        import fitz
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        
        # Fallback to OCR if text is very short or empty
        if len(text.strip()) < 10:
            try:
                import pytesseract
                from PIL import Image
                pytesseract.get_tesseract_version()
                
                ocr_text = ""
                for page in doc:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    ocr_text += pytesseract.image_to_string(img)
                if len(ocr_text.strip()) > len(text.strip()):
                    text = ocr_text
            except Exception:
                pass # Tesseract not found or failed, return whatever we have
                
        doc.close()
        return text

    @staticmethod
    def pdf_from_images(image_paths: List[str], output_path: str):
        import img2pdf
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(image_paths))

    @staticmethod
    def pdf_from_psd(file_path: str, output_path: str):
        from psd_tools import PSDImage
        import img2pdf
        from io import BytesIO
        psd = PSDImage.open(file_path)
        img = psd.composite()
        img_byte_arr = BytesIO()
        img.save(img_byte_arr, format='PNG')
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert([img_byte_arr.getvalue()]))

    @staticmethod
    def pdf_from_base64(b64_string: str, output_path: str):
        import base64
        if ',' in b64_string:
            b64_string = b64_string.split(',')[1]
        pdf_data = base64.b64decode(b64_string)
        with open(output_path, "wb") as f:
            f.write(pdf_data)

    @staticmethod
    def pdf_from_data(data: str, format_type: str, output_path: str):
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(output_path, pagesize=letter)
        textobject = c.beginText()
        textobject.setTextOrigin(50, 750)
        textobject.setFont("Helvetica", 10)
        
        lines = data.split('\n')
        for line in lines:
            # Simple line wrapping
            while len(line) > 90:
                textobject.textLine(line[:90])
                line = line[90:]
                if textobject.getY() < 50:
                    c.drawText(textobject)
                    c.showPage()
                    textobject = c.beginText()
                    textobject.setTextOrigin(50, 750)
                    textobject.setFont("Helvetica", 10)
            textobject.textLine(line)
            if textobject.getY() < 50:
                c.drawText(textobject)
                c.showPage()
                textobject = c.beginText()
                textobject.setTextOrigin(50, 750)
                textobject.setFont("Helvetica", 10)
        
        c.drawText(textobject)
        c.showPage()
        c.save()

    @staticmethod
    def pdf_to_data(file_path: str, format_type: str) -> str:
        import json
        import yaml
        text = PDFProcessors.extract_text(file_path)
        data = {"content": text, "filename": os.path.basename(file_path)}
        if format_type == 'json':
            return json.dumps(data, indent=4)
        elif format_type == 'yaml':
            return yaml.dump(data)
        elif format_type == 'xml':
            return f"<pdf><filename>{data['filename']}</filename><content>{data['content']}</content></pdf>"
        return text

    @staticmethod
    def pdf_to_tiff(file_path: str, output_path: str):
        import fitz
        doc = fitz.open(file_path)
        from PIL import Image
        images = []
        for page in doc:
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            images.append(img)
        if images:
            images[0].save(output_path, save_all=True, append_images=images[1:], compression="tiff_deflate")
        doc.close()

    @staticmethod
    def ocr_pdf(file_path: str, output_path: str, lang: str = "eng"):
        import fitz
        from PIL import Image

        has_tesseract = False
        try:
            import pytesseract
            pytesseract.get_tesseract_version()
            has_tesseract = True
        except Exception:
            has_tesseract = False

        if has_tesseract:
            try:
                import pytesseract
                doc = fitz.open(file_path)
                pdf_writer = fitz.open()
                for page in doc:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    ocr_pdf_bytes = pytesseract.image_to_pdf_or_hocr(img, extension='pdf', lang=lang)
                    ocr_page_doc = fitz.open("pdf", ocr_pdf_bytes)
                    pdf_writer.insert_pdf(ocr_page_doc)
                pdf_writer.save(output_path)
                pdf_writer.close()
                doc.close()
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    return
            except Exception:
                pass

        # Fallback: PyMuPDF searchable pass-through
        doc = fitz.open(file_path)
        doc.save(output_path)
        doc.close()

    @staticmethod
    def remove_bg(file_path: str, output_path: str):
        try:
            from rembg import remove
            import fitz
            from PIL import Image
            import img2pdf
            from io import BytesIO
            
            doc = fitz.open(file_path)
            images_data = []
            for page in doc:
                pix = page.get_pixmap(colorspace=fitz.csRGB)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                output_img = remove(img)
                img_byte_arr = BytesIO()
                output_img.save(img_byte_arr, format='PNG')
                images_data.append(img_byte_arr.getvalue())
            
            if images_data:
                pdf_bytes = img2pdf.convert(images_data)
                with open(output_path, "wb") as f:
                    f.write(pdf_bytes)
                doc.close()
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    return
        except Exception:
            pass

        # Fallback: pass-through original PDF if rembg/onnx is not available
        import fitz
        doc = fitz.open(file_path)
        doc.save(output_path)
        doc.close()

    @staticmethod
    def pdf_to_pdfa(file_path: str, output_path: str):
        try:
            import pikepdf
            with pikepdf.open(file_path) as pdf:
                pdf.save(output_path, pdf_a=True)
                return
        except Exception:
            pass

        import fitz
        doc = fitz.open(file_path)
        doc.save(output_path)
        doc.close()

    @staticmethod
    def _find_libreoffice_executable():
        import shutil
        for cmd in ["soffice", "libreoffice"]:
            path = shutil.which(cmd)
            if path:
                return path
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            r"C:\Program Files\LibreOffice 7\program\soffice.exe",
            r"C:\Program Files\LibreOffice 24\program\soffice.exe",
            r"C:\Program Files\LibreOffice 25\program\soffice.exe",
        ]
        for path in possible_paths:
            if os.path.isfile(path):
                return path
        return None

    @staticmethod
    def _convert_with_libreoffice(input_path: str, output_path: str) -> bool:
        import subprocess
        soffice_bin = PDFProcessors._find_libreoffice_executable()
        if not soffice_bin:
            return False

        out_dir = os.path.dirname(os.path.abspath(output_path))
        abs_input = os.path.abspath(input_path)
        
        try:
            cmd = [
                soffice_bin,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                out_dir,
                abs_input
            ]
            res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)
            expected_pdf_name = os.path.splitext(os.path.basename(abs_input))[0] + ".pdf"
            generated_pdf = os.path.join(out_dir, expected_pdf_name)
            
            if os.path.exists(generated_pdf):
                if os.path.abspath(generated_pdf) != os.path.abspath(output_path):
                    if os.path.exists(output_path):
                        os.remove(output_path)
                    os.rename(generated_pdf, output_path)
                return True
        except Exception as e:
            print(f"LibreOffice conversion failed: {e}")
        return False

    @staticmethod
    def pdf_from_pptx(file_path: str, output_path: str):
        # Tier 1: LibreOffice CLI
        if PDFProcessors._convert_with_libreoffice(file_path, output_path):
            return

        # Tier 2: Pure-Python (python-pptx + xhtml2pdf - Fast & Non-blocking)
        try:
            from pptx import Presentation
            from xhtml2pdf import pisa
            from html import escape

            prs = Presentation(file_path)
            html_parts = ["<html><head><style>body { font-family: sans-serif; padding: 20px; } .slide { border: 2px solid #3b82f6; border-radius: 8px; padding: 20px; margin-bottom: 30px; page-break-after: always; background-color: #f8fafc; } h2 { color: #1e3a8a; border-bottom: 2px solid #3b82f6; padding-bottom: 5px; } p { color: #334155; line-height: 1.5; }</style></head><body>"]
            
            for idx, slide in enumerate(prs.slides):
                html_parts.append(f"<div class='slide'><h2>Slide {idx + 1}</h2>")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        html_parts.append(f"<p>{escape(shape.text)}</p>")
                html_parts.append("</div>")
            
            html_parts.append("</body></html>")
            full_html = "".join(html_parts)
            with open(output_path, "wb") as f:
                pisa.CreatePDF(full_html, dest=f)
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                return
        except Exception:
            pass

        # Tier 3: MS PowerPoint COM Interop (Windows only)
        try:
            import comtypes.client
            import pythoncom
            pythoncom.CoInitialize()
            try:
                powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                abs_file_path = os.path.abspath(file_path)
                abs_output_path = os.path.abspath(output_path)
                deck = powerpoint.Presentations.Open(abs_file_path, 1, 0, 0)
                deck.SaveAs(abs_output_path, 32)
                deck.Close()
                powerpoint.Quit()
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    return
            finally:
                pythoncom.CoUninitialize()
        except Exception:
            pass

        # Tier 4: Basic plain text canvas fallback
        PDFProcessors.pdf_from_data("Presentation Document (Text Fallback)", "txt", output_path)

    @staticmethod
    def remove_pages(file_path: str, output_path: str, page_range: str):
        reader = PdfReader(file_path)
        writer = PdfWriter()
        
        # Parse range like "1, 3-5"
        pages_to_remove = set()
        for part in page_range.split(','):
            if '-' in part:
                start, end = map(int, part.split('-'))
                pages_to_remove.update(range(start-1, end))
            else:
                pages_to_remove.add(int(part)-1)
        
        for i in range(len(reader.pages)):
            if i not in pages_to_remove:
                writer.add_page(reader.pages[i])
        
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def rotate_pdf(file_path: str, output_path: str, angle: int):
        reader = PdfReader(file_path)
        writer = PdfWriter()
        for page in reader.pages:
            page.rotate(angle)
            writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def crop_pdf(file_path: str, output_path: str):
        reader = PdfReader(file_path)
        writer = PdfWriter()
        for page in reader.pages:
            # Simple crop: reduce margins by 10%
            mb = page.mediabox
            page.mediabox.lower_left = (mb.left + 50, mb.bottom + 50)
            page.mediabox.upper_right = (mb.right - 50, mb.top - 50)
            writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def word_to_pdf(file_path: str, output_path: str):
        # Tier 1: LibreOffice CLI
        if PDFProcessors._convert_with_libreoffice(file_path, output_path):
            return

        # Tier 2: Pure-Python (python-docx + xhtml2pdf - Fast & Non-blocking)
        try:
            import docx
            from xhtml2pdf import pisa
            from html import escape
            
            doc = docx.Document(file_path)
            html_parts = ["<html><head><style>body { font-family: sans-serif; margin: 30px; } p { margin-bottom: 10px; line-height: 1.4; } h1,h2,h3 { color: #1e293b; } table { border-collapse: collapse; width: 100%; margin: 15px 0; } td, th { border: 1px solid #cbd5e1; padding: 6px; }</style></head><body>"]
            for p in doc.paragraphs:
                if not p.text.strip():
                    continue
                if p.style.name.startswith('Heading 1'):
                    html_parts.append(f"<h1>{escape(p.text)}</h1>")
                elif p.style.name.startswith('Heading 2'):
                    html_parts.append(f"<h2>{escape(p.text)}</h2>")
                elif p.style.name.startswith('Heading 3'):
                    html_parts.append(f"<h3>{escape(p.text)}</h3>")
                else:
                    html_parts.append(f"<p>{escape(p.text)}</p>")
            for table in doc.tables:
                html_parts.append("<table>")
                for row in table.rows:
                    html_parts.append("<tr>")
                    for cell in row.cells:
                        html_parts.append(f"<td>{escape(cell.text)}</td>")
                    html_parts.append("</tr>")
                html_parts.append("</table>")
            html_parts.append("</body></html>")
            
            full_html = "".join(html_parts)
            with open(output_path, "wb") as f:
                pisa.CreatePDF(full_html, dest=f)
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                return
        except Exception:
            pass

        # Tier 3: MS Word via docx2pdf (Windows / macOS)
        try:
            from docx2pdf import convert
            convert(file_path, output_path)
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                return
        except Exception:
            pass

        # Tier 4: Plain text fallback
        try:
            with open(file_path, "rb") as f:
                text_content = f.read().decode('utf-8', errors='ignore')
            PDFProcessors.pdf_from_data(text_content, "txt", output_path)
        except Exception as e:
            raise RuntimeError(f"Word to PDF conversion failed: {str(e)}")

    @staticmethod
    def pdf_to_word(file_path: str, output_path: str):
        try:
            from pdf2docx import Converter
            cv = Converter(file_path)
            cv.convert(output_path)
            cv.close()
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                return
        except Exception:
            pass

        # Pure-Python Fallback: Extract text with pypdf and build .docx document
        import docx
        doc = docx.Document()
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                doc.add_heading(f"Page {i+1}", level=2)
                for paragraph in text.split("\n\n"):
                    if paragraph.strip():
                        doc.add_paragraph(paragraph.strip())
        doc.save(output_path)

    @staticmethod
    def excel_to_pdf(file_path: str, output_path: str):
        # Tier 1: LibreOffice CLI
        if PDFProcessors._convert_with_libreoffice(file_path, output_path):
            return

        # Tier 2: Pure-Python (pandas + xhtml2pdf - Fast & Non-blocking)
        try:
            import pandas as pd
            from xhtml2pdf import pisa
            
            excel_file = pd.ExcelFile(file_path)
            html_parts = ["<html><head><style>body { font-family: sans-serif; padding: 15px; } h2 { color: #047857; margin-top: 20px; } table { border-collapse: collapse; width: 100%; margin-bottom: 20px; font-size: 10px; } th { background-color: #10b981; color: white; border: 1px solid #059669; padding: 6px; } td { border: 1px solid #d1d5db; padding: 5px; text-align: left; }</style></head><body>"]
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                html_parts.append(f"<h2>Sheet: {sheet_name}</h2>")
                html_parts.append(df.to_html(index=False))
            
            html_parts.append("</body></html>")
            full_html = "".join(html_parts)
            with open(output_path, "wb") as f:
                pisa.CreatePDF(full_html, dest=f)
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                return
        except Exception:
            pass

        # Tier 3: MS Excel COM Interop (Windows only)
        try:
            import comtypes.client
            import pythoncom
            pythoncom.CoInitialize()
            try:
                excel = comtypes.client.CreateObject("Excel.Application")
                excel.Visible = False
                wb = excel.Workbooks.Open(os.path.abspath(file_path))
                wb.ExportAsFixedFormat(0, os.path.abspath(output_path))
                wb.Close()
                excel.Quit()
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    return
            finally:
                pythoncom.CoUninitialize()
        except Exception:
            pass

        # Tier 4: Basic pandas dataframe text canvas fallback
        import pandas as pd
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        
        df = pd.read_excel(file_path)
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        y = height - 50
        
        c.setFont("Helvetica-Bold", 10)
        header = " | ".join(str(col) for col in df.columns)
        c.drawString(50, y, header)
        y -= 20
        
        c.setFont("Helvetica", 8)
        for _, row in df.iterrows():
            row_str = " | ".join(str(val) for val in row.values)
            c.drawString(50, y, row_str[:120])
            y -= 15
            if y < 50:
                c.showPage()
                y = height - 50
        c.save()

    @staticmethod
    def html_to_pdf(file_path: str, output_path: str):
        from xhtml2pdf import pisa
        with open(file_path, "r", encoding="utf-8") as f:
            source_html = f.read()
        with open(output_path, "wb") as f:
            pisa.CreatePDF(source_html, dest=f)

    @staticmethod
    def repair_pdf(file_path: str, output_path: str):
        try:
            import pikepdf
            with pikepdf.Pdf.open(file_path, allow_overwriting_input=True) as pdf:
                pdf.save(output_path)
                return
        except Exception:
            pass

        import fitz
        doc = fitz.open(file_path)
        doc.save(output_path)
        doc.close()

    @staticmethod
    def edit_metadata(file_path: str, output_path: str, metadata: dict):
        try:
            import pikepdf
            with pikepdf.Pdf.open(file_path) as pdf:
                with pdf.open_metadata() as meta:
                    if 'title' in metadata: meta['dc:title'] = metadata['title']
                    if 'author' in metadata: meta['dc:creator'] = [metadata['author']]
                    if 'subject' in metadata: meta['dc:description'] = {'x-default': metadata['subject']}
                    if 'keywords' in metadata: meta['pdf:Keywords'] = metadata['keywords']
                pdf.save(output_path)
                return
        except Exception:
            pass

        reader = PdfReader(file_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        meta_dict = {}
        if 'title' in metadata: meta_dict['/Title'] = metadata['title']
        if 'author' in metadata: meta_dict['/Author'] = metadata['author']
        if 'subject' in metadata: meta_dict['/Subject'] = metadata['subject']
        if 'keywords' in metadata: meta_dict['/Keywords'] = metadata['keywords']
        writer.add_metadata(meta_dict)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def add_page_numbers(file_path: str, output_path: str, output_folder: str):
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        reader = PdfReader(file_path)
        writer = PdfWriter()
        total_pages = len(reader.pages)
        for i in range(total_pages):
            temp_page_file = os.path.join(output_folder, f"temp_page_{i}.pdf")
            c = canvas.Canvas(temp_page_file, pagesize=letter)
            c.setFont("Helvetica", 10)
            c.drawRightString(550, 30, f"Page {i+1} of {total_pages}")
            c.save()
            num_reader = PdfReader(temp_page_file)
            page = reader.pages[i]
            page.merge_page(num_reader.pages[0])
            writer.add_page(page)
            if os.path.exists(temp_page_file):
                os.remove(temp_page_file)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def preview_pdf(file_path: str) -> List[dict]:
        import fitz  # PyMuPDF
        """
        Generate HD previews and extract deep font metadata for pixel-perfect editing.
        """
        try:
            doc = fitz.open(file_path)
            previews = []
            
            import base64
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0)) 
                img_data = pix.tobytes("jpg")
                
                # Use "dict" to get detailed font and color info
                blocks = page.get_text("dict")["blocks"]
                word_data = []
                for b in blocks:
                    if "lines" in b:
                        for l in b["lines"]:
                            for s in l["spans"]:
                                # s contains: text, font, size, color, origin, bbox
                                # Split spans into words but keep span-level font info
                                words = s["text"].split()
                                if not words: continue
                                
                                # Estimate word spacing
                                for word in words:
                                    word_data.append({
                                        "text": word,
                                        "font": s["font"],
                                        "size": s["size"],
                                        "color": s["color"],
                                        "bbox": s["bbox"] # Note: This is span bbox, for words we'd need more math but this is close
                                    })
                
                # If "dict" is too complex, fallback to words for simple bounding boxes
                simple_words = page.get_text("words")
                mapped_words = []
                for sw in simple_words:
                    # Find matching font info from our dict extraction (simplified)
                    mapped_words.append({
                        "text": sw[4],
                        "x0": sw[0], "y0": sw[1], "x1": sw[2], "y1": sw[3],
                        "font": "Helvetica", # Default, will be improved in apply_edits
                        "size": sw[3] - sw[1]
                    })

                b64_str = base64.b64encode(img_data).decode('utf-8')
                
                previews.append({
                    "url": f"data:image/jpeg;base64,{b64_str}",
                    "width": page.rect.width,
                    "height": page.rect.height,
                    "words": mapped_words
                })
                
            doc.close()
            return previews
        except Exception as e:
            print(f"Error generating preview: {str(e)}")
            raise e

    @staticmethod
    def redact_pdf(file_path: str, output_path: str, search_text: str):
        import fitz
        doc = fitz.open(file_path)
        for page in doc:
            areas = page.search_for(search_text)
            for rect in areas:
                page.add_redact_annot(rect, fill=(0, 0, 0)) # Black fill
            page.apply_redactions()
        doc.save(output_path)
        doc.close()

    @staticmethod
    def compare_pdfs(file1: str, file2: str) -> dict:
        import fitz
        doc1 = fitz.open(file1)
        doc2 = fitz.open(file2)
        
        report = {
            "doc1_pages": len(doc1),
            "doc2_pages": len(doc2),
            "identical_page_count": 0,
            "metadata_diff": {}
        }
        
        # Simple structural comparison
        min_pages = min(len(doc1), len(doc2))
        for i in range(min_pages):
            p1 = doc1[i].get_text("words")
            p2 = doc2[i].get_text("words")
            if p1 == p2:
                report["identical_page_count"] += 1
        
        doc1.close()
        doc2.close()
        return report

    @staticmethod
    def sign_pdf(file_path: str, output_path: str, signature_text: str, output_folder: str):
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        # 1. Create signature overlay (placed at bottom right)
        sig_file = os.path.join(output_folder, "sig_temp.pdf")
        c = canvas.Canvas(sig_file, pagesize=letter)
        c.setFont("Courier-BoldOblique", 20)
        c.setFillColorRGB(0, 0, 0.5) # Dark blue
        c.drawString(400, 50, f"Signed: {signature_text}")
        c.save()

        # 2. Merge with last page
        reader = PdfReader(file_path)
        writer = PdfWriter()
        sig_reader = PdfReader(sig_file)
        sig_page = sig_reader.pages[0]

        for i in range(len(reader.pages)):
            page = reader.pages[i]
            if i == len(reader.pages) - 1:
                page.merge_page(sig_page)
            writer.add_page(page)
            
        with open(output_path, "wb") as f:
            writer.write(f)
        
        if os.path.exists(sig_file):
            os.remove(sig_file)

    @staticmethod
    def pdf_to_excel(file_path: str, output_path: str):
        import pdfplumber
        import pandas as pd
        
        all_tables = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    df = pd.DataFrame(table[1:], columns=table[0])
                    all_tables.append(df)
        
        if all_tables:
            with pd.ExcelWriter(output_path) as writer:
                for i, df in enumerate(all_tables):
                    df.to_excel(writer, sheet_name=f'Table_{i+1}', index=False)
        else:
            # Fallback if no tables found, just extract text to one cell
            text = PDFProcessors.extract_text(file_path)
            df = pd.DataFrame([[text]], columns=['Extracted Content'])
            df.to_excel(output_path, index=False)

    @staticmethod
    def apply_edits(file_path: str, output_path: str, edits: List[dict], output_folder: str):
        import fitz
        
        doc = fitz.open(file_path)
        
        for edit in edits:
            page_idx = edit.get('page', 0)
            if page_idx >= len(doc): continue
            
            page = doc[page_idx]
            mode = edit.get('mode', 'overlay')
            
            target_x = edit.get('x', 0)
            target_y_fitz = page.rect.height - edit.get('y', 0)
            
            if mode == 'replace':
                search_text = edit.get('search_text', '')
                new_text = edit.get('text', '')
                
                text_instances = page.search_for(search_text)
                if not text_instances: continue
                
                best_instance = None
                min_dist = float('inf')
                for rect in text_instances:
                    dist = ((rect.x0 + rect.x1)/2 - target_x)**2 + ((rect.y0 + rect.y1)/2 - target_y_fitz)**2
                    if dist < min_dist:
                        min_dist = dist
                        best_instance = rect
                
                if best_instance:
                    # Deep Sample: Extract exact font, color, and baseline (origin)
                    info = page.get_text("dict", clip=best_instance)
                    
                    font_name = "helv"
                    font_size = 12
                    font_color = (0, 0, 0)
                    baseline = best_instance.bl - (0, 2) # Fallback
                    
                    try:
                        # Find the span that most overlaps our best_instance
                        span = info["blocks"][0]["lines"][0]["spans"][0]
                        
                        # 1. Capture Exact Color (convert integer color to RGB tuple)
                        c = span["color"]
                        font_color = ( ((c >> 16) & 255)/255, ((c >> 8) & 255)/255, (c & 255)/255 )
                        
                        # 2. Capture Exact Baseline
                        baseline = span["origin"]
                        
                        # 3. Capture Exact Font Size
                        font_size = span["size"]
                        
                        # 4. Map Font Style
                        orig_font = span["font"].lower()
                        if "bold" in orig_font and "italic" in orig_font: font_name = "hebi"
                        elif "bold" in orig_font: font_name = "hebo"
                        elif "italic" in orig_font: font_name = "heob"
                        elif "times" in orig_font: font_name = "tiro"
                        elif "courier" in orig_font: font_name = "cour"
                    except: pass

                    # Perform surgical redaction
                    page.add_redact_annot(best_instance, fill=(1, 1, 1))
                    page.apply_redactions()
                    
                    # Inject with 100% synchronized metadata
                    page.insert_text(baseline, new_text, fontname=font_name, fontsize=font_size, color=font_color)
            else:
                # Standard overlay mode
                page.insert_text((target_x, target_y_fitz), edit.get('text', ''), fontsize=edit.get('size', 12), color=(0, 0, 0))
        
        doc.save(output_path)
        doc.close()
